# Standard library imports
import io
import json
import mimetypes
import os
from typing import Tuple, List, Dict

# Third-party imports
from dotenv import load_dotenv
from flask import Flask, request, jsonify
from PIL import Image, ExifTags
from bs4 import BeautifulSoup

# PDF / OCR / Tabelas
import pdfplumber
from pdf2image import convert_from_bytes
import pytesseract
import camelot

# Office / Dados
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

# STT opcional (offline)
try:
    from vosk import Model, KaldiRecognizer
    VOSK_AVAILABLE = True
except Exception:
    VOSK_AVAILABLE = False

# Define supported file types and their MIME types
SUPPORTED_FORMATS = {
    "pdf": ["application/pdf"],
    "powerpoint": [
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ],
    "word": [
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ],
    "excel": [
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "text/csv",
    ],
    "image": [
        "image/jpeg",
        "image/png",
        "image/gif",
        "image/bmp",
        "image/tiff",
        "image/webp",
    ],
    "audio": [
        "audio/mpeg",
        "audio/wav",
        "audio/ogg",
        "audio/m4a",
        "audio/mp3",
        "audio/aac",
    ],
    "html": ["text/html"],
    "text": [
        "text/plain",
        "application/json",
        "application/xml",
        "text/xml",
    ],
}

# Initialize Flask app
app = Flask(__name__)


def is_supported_format(content_type: str) -> Tuple[bool, str]:
    """
    Check if the content type is supported and return format type
    """
    if not content_type:
        return False, ""
    for format_type, mime_types in SUPPORTED_FORMATS.items():
        if any(content_type.lower().startswith(mime) for mime in mime_types):
            return True, format_type
    return False, ""


# ---------------------
# PDF handlers (texto + tabelas + OCR opcional)
# ---------------------
def pdf_extract_text_and_tables(pdf_bytes: bytes, use_ocr: bool) -> Dict:
    pages_text: List[str] = []
    tables: List[Dict] = []
    used_ocr = False

    # 1) Tenta texto nativo
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                txt = page.extract_text() or ""
                pages_text.append(txt)
    except Exception:
        pages_text = []

    # 2) Extrai tabelas com Camelot (precisa arquivo temporário)
    try:
        tmp_dir = "uploads"
        os.makedirs(tmp_dir, exist_ok=True)
        tmp_path = os.path.join(tmp_dir, "temp_pdf.pdf")
        with open(tmp_path, "wb") as f:
            f.write(pdf_bytes)

        # "stream" funciona bem para muitas tabelas; "lattice" para tabelas com bordas
        c_tables = camelot.read_pdf(tmp_path, pages="all", flavor="stream")
        for t in c_tables:
            tables.append({
                "page": t.page,
                "rows": t.df.fillna("").values.tolist()
            })

        try:
            os.remove(tmp_path)
        except Exception:
            pass
    except Exception:
        tables = []

    # 3) Decide se precisa de OCR
    def looks_like_image(pages_txt: List[str]) -> bool:
        if not pages_txt:
            return True
        empties = sum(1 for t in pages_txt if (t or "").strip() == "")
        return empties >= max(1, len(pages_txt) // 2)

    if use_ocr or looks_like_image(pages_text):
        try:
            images = convert_from_bytes(pdf_bytes, dpi=300)
            pages_text = []
            lang = os.getenv("OCR_LANG", "por+eng")
            for img in images:
                txt = pytesseract.image_to_string(img, lang=lang)
                pages_text.append(txt or "")
            used_ocr = True
            tables = []  # via OCR perde-se a estrutura de tabela
        except Exception:
            pass

    content = "\n\n".join(pages_text).strip()
    return {"content": content, "pages": pages_text, "tables": tables, "ocr": used_ocr}


# ---------------------
# Word / PowerPoint / Excel / Imagens / HTML / Texto / Áudio
# ---------------------
def extract_word(path: str) -> str:
    doc = DocxDocument(path)
    parts = []
    for p in doc.paragraphs:
        parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join([cell.text for cell in row.cells]))
    return "\n".join([p for p in parts if p and p.strip()])


def extract_powerpoint(path: str) -> str:
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def extract_excel(path: str, content_type: str) -> Dict:
    out_tables = []
    content_parts = []
    if content_type.endswith("csv"):
        df = pd.read_csv(path, dtype=str, keep_default_na=False)
        out_tables.append({"sheet": "CSV", "rows": [df.columns.tolist()] + df.values.tolist()})
        content_parts.append(df.to_csv(index=False))
    else:
        xls = pd.ExcelFile(path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name, dtype=str, keep_default_na=False)
            out_tables.append({"sheet": sheet_name, "rows": [df.columns.tolist()] + df.values.tolist()})
            content_parts.append(f"## {sheet_name}\n" + df.to_csv(index=False))
    return {"content": "\n\n".join(content_parts), "tables": out_tables}


def extract_image(path: str) -> Dict:
    img = Image.open(path)
    # EXIF
    exif_data = {}
    try:
        exif = img._getexif() or {}
        for tag_id, val in exif.items():
            tag = ExifTags.TAGS.get(tag_id, tag_id)
            exif_data[tag] = str(val)
    except Exception:
        exif_data = {}
    # OCR
    try:
        text = pytesseract.image_to_string(img, lang=os.getenv("OCR_LANG", "por+eng"))
    except Exception:
        text = ""
    return {"content": text.strip(), "exif": exif_data}


def extract_html(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        html = f.read()
    soup = BeautifulSoup(html, "lxml")
    for script in soup(["script", "style"]):
        script.extract()
    text = soup.get_text(separator="\n")
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    return "\n".join(lines)


def extract_text_generic(path: str, content_type: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = f.read()
        if content_type.endswith("json"):
            try:
                return json.dumps(json.loads(data), ensure_ascii=False, indent=2)
            except Exception:
                return data
        return data
    except Exception:
        return ""


def transcribe_audio(path: str) -> Dict:
    if not VOSK_AVAILABLE:
        return {"error": "Audio transcription disabled: Vosk not installed/available."}
    model_path = os.getenv("VOSK_MODEL_PATH", "")
    if not model_path or not os.path.isdir(model_path):
        return {"error": "Set VOSK_MODEL_PATH to a valid local model directory."}

    import wave
    import subprocess

    # Garante WAV mono 16kHz (converte via ffmpeg se necessário)
    wav_path = path
    try:
        with wave.open(path, "rb") as w:
            params_ok = (w.getnchannels() == 1 and w.getframerate() in (16000, 8000))
    except Exception:
        params_ok = False
    if not params_ok:
        wav_path = path + ".conv.wav"
        subprocess.run(
            ["ffmpeg", "-y", "-i", path, "-ac", "1", "-ar", "16000", wav_path],
            check=False,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )

    wf = wave.open(wav_path, "rb")
    model = Model(model_path)
    rec = KaldiRecognizer(model, wf.getframerate())
    rec.SetWords(True)
    result_text = []
    while True:
        data = wf.readframes(4000)
        if len(data) == 0:
            break
        if rec.AcceptWaveform(data):
            res = json.loads(rec.Result())
            if "text" in res:
                result_text.append(res["text"])
    final = json.loads(rec.FinalResult()).get("text", "")
    if final:
        result_text.append(final)

    if wav_path.endswith(".conv.wav"):
        try:
            os.remove(wav_path)
        except Exception:
            pass

    return {"content": " ".join(result_text).strip()}


@app.route("/health", methods=["GET"])
def health():
    """Health check endpoint"""
    return jsonify({"status": "healthy"}), 200


@app.route("/convert", methods=["POST"])
def convert():
    try:
        file_data = None
        content_type = None

        # 1) Se vier como multipart/form-data (ex.: n8n mandando Binary como 'file')
        if "file" in request.files:
            f = request.files["file"]
            file_data = f.read()
            content_type = f.mimetype or request.content_type

        # 2) Caso contrário, RAW no corpo
        if file_data is None:
            file_data = request.get_data()
            content_type = request.content_type or ""

        if not file_data:
            return jsonify({"error": "No file data provided"}), 400

        # Determine file extension from content type
        is_supported, format_type = is_supported_format(content_type)
        if not is_supported:
            return jsonify({"error": "Unsupported Media Type"}), 415

        extension = mimetypes.guess_extension(content_type) or ""
        if extension == ".jpe":  # quirk do mimetypes para image/jpeg
            extension = ".jpg"
        temp_filename = f"temp_file{extension}"
        temp_path = os.path.join("uploads", temp_filename)

        # Save the binary data to a temporary file
        os.makedirs("uploads", exist_ok=True)
        with open(temp_path, "wb") as f:
            f.write(file_data)

        try:
            use_ocr = request.args.get("ocr", "true").lower() == "true"

            if format_type == "pdf":
                app.logger.info("Processing PDF file...")
                result = pdf_extract_text_and_tables(file_data, use_ocr=use_ocr)
                app.logger.info(f"PDF processing result: {result}")
                return jsonify({"format": "pdf", **result})

            if format_type == "image":
                result = extract_image(temp_path)
                # Para imagem, OCR é sempre relevante
                result["ocr"] = True
                return jsonify({"format": "image", **result})

            if format_type == "word":
                content = extract_word(temp_path)
                return jsonify({"format": "word", "content": content})

            if format_type == "powerpoint":
                content = extract_powerpoint(temp_path)
                return jsonify({"format": "powerpoint", "content": content})

            if format_type == "excel":
                result = extract_excel(temp_path, content_type)
                return jsonify({"format": "excel", **result})

            if format_type == "html":
                content = extract_html(temp_path)
                return jsonify({"format": "html", "content": content})

            if format_type == "text":
                content = extract_text_generic(temp_path, content_type)
                return jsonify({"format": "text", "content": content})

            if format_type == "audio":
                result = transcribe_audio(temp_path)
                return jsonify({"format": "audio", **result})

            return jsonify({"error": "Unhandled format"}), 400

        finally:
            if os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except Exception:
                    pass

    except Exception as e:
        app.logger.error(f"Error during conversion: {e}")
        return jsonify({"error": str(e)}), 500


if __name__ == "__main__":
    # Load environment variables
    load_dotenv()
    # Create uploads directory if it doesn't exist
    os.makedirs("uploads", exist_ok=True)
    # Run the Flask app
    app.run(host="0.0.0.0", port=7001)
